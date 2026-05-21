"""
VisionWise Capital — template builder.

Regenerates all brand-aligned DOCX/PPTX templates in templates/ from
the official brand spec extracted into company-intel/brand/.

Run from repo root:  python company-intel/brand/build-templates.py
"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Cm, Inches, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches as PInches, Pt as PPt

ROOT = Path(__file__).resolve().parent.parent.parent
LOGO_PNG = ROOT / "company-intel/brand/logo/logo-rgb.png"
LOGO_SMALL = ROOT / "company-intel/brand/logo/logo-small.png"
ICON_PNG = ROOT / "company-intel/brand/logo/icon.png"
TEMPLATES = ROOT / "templates"
TEMPLATES.mkdir(exist_ok=True)

GREEN = RGBColor(0xA3, 0xD5, 0x5D)
GREEN_DARK = RGBColor(0x5C, 0x96, 0x27)
GREEN_TINT = RGBColor(0xDA, 0xF0, 0xBC)
CHARCOAL = RGBColor(0x40, 0x46, 0x4B)
MID = RGBColor(0x8A, 0x8F, 0x94)
LIGHT = RGBColor(0xE5, 0xE7, 0xEA)
OFFWHITE = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

P_GREEN = PRGB(0xA3, 0xD5, 0x5D)
P_GREEN_DARK = PRGB(0x5C, 0x96, 0x27)
P_GREEN_TINT = PRGB(0xDA, 0xF0, 0xBC)
P_CHARCOAL = PRGB(0x40, 0x46, 0x4B)
P_MID = PRGB(0x8A, 0x8F, 0x94)
P_LIGHT = PRGB(0xE5, 0xE7, 0xEA)
P_WHITE = PRGB(0xFF, 0xFF, 0xFF)

# VWC holds licenses for Gill Sans MT Pro Book and Proxima Nova Regular,
# so the templates specify them directly. Word/PowerPoint will substitute
# on machines that lack the fonts; documented Google Fonts fallback chain:
#   Gill Sans MT Pro Book  →  Lato (Google Fonts)   →  Helvetica / Arial
#   Proxima Nova Regular   →  Open Sans (Google Fonts)  →  Helvetica / Arial
HEADING_FONT = "Gill Sans MT Pro"
BODY_FONT = "Proxima Nova"

DISCLAIMER = (
    "This communication is provided for informational purposes only and does not "
    "constitute an offer to sell or a solicitation of an offer to buy any securities. "
    "Any such offer or solicitation will be made only through definitive offering "
    "documents (including a private placement memorandum) provided to qualified "
    "investors. Past performance is not indicative of future results. Real estate "
    "investments involve risk, including loss of principal."
)


# ---------------------------------------------------------------------------
# DOCX helpers
# ---------------------------------------------------------------------------

def set_font(run, name=BODY_FONT, size=11, bold=False, color=CHARCOAL, italic=False):
    run.font.name = name
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    for attr in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
        rfonts.set(qn(attr), name)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_para(doc_or_cell, text, font=BODY_FONT, size=11, bold=False, color=CHARCOAL,
             align=None, space_after=6, italic=False):
    if hasattr(doc_or_cell, "add_paragraph"):
        p = doc_or_cell.add_paragraph()
    else:
        p = doc_or_cell.paragraphs[0] if not doc_or_cell.text else doc_or_cell.add_paragraph()
    if align is not None:
        p.alignment = align
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.space_before = Pt(0)
    if text:
        run = p.add_run(text)
        set_font(run, name=font, size=size, bold=bold, color=color, italic=italic)
    return p


def h1(doc, text):
    return add_para(doc, text, font=HEADING_FONT, size=28, bold=True, color=CHARCOAL, space_after=12)


def h2(doc, text):
    p = add_para(doc, text, font=HEADING_FONT, size=18, bold=True, color=CHARCOAL, space_after=4)
    add_bottom_border(p, GREEN, size_pt=2)
    p.paragraph_format.space_after = Pt(10)
    return p


def h3(doc, text):
    return add_para(doc, text, font=HEADING_FONT, size=13, bold=True, color=GREEN_DARK, space_after=4)


def caption(doc, text):
    return add_para(doc, text, font=BODY_FONT, size=9, color=MID, space_after=4, italic=True)


def add_bottom_border(p, color: RGBColor, size_pt=1):
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), str(int(size_pt * 8)))
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    pBdr.append(bottom)
    pPr.append(pBdr)


def shade_cell(cell, color: RGBColor):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    tcPr.append(shd)


def set_cell_borders(cell, color: RGBColor = LIGHT, size_pt=0.5):
    tcPr = cell._tc.get_or_add_tcPr()
    tcBorders = OxmlElement("w:tcBorders")
    for edge in ("top", "left", "bottom", "right"):
        b = OxmlElement(f"w:{edge}")
        b.set(qn("w:val"), "single")
        b.set(qn("w:sz"), str(int(size_pt * 8)))
        b.set(qn("w:space"), "0")
        b.set(qn("w:color"), f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
        tcBorders.append(b)
    tcPr.append(tcBorders)


def fill_cell(cell, text, font=BODY_FONT, size=10, bold=False, color=CHARCOAL,
              fill: RGBColor | None = None, align=None, vert="center"):
    cell.vertical_alignment = (
        WD_ALIGN_VERTICAL.CENTER if vert == "center"
        else WD_ALIGN_VERTICAL.TOP if vert == "top"
        else WD_ALIGN_VERTICAL.BOTTOM
    )
    p = cell.paragraphs[0]
    p.text = ""
    if align is not None:
        p.alignment = align
    p.paragraph_format.space_after = Pt(0)
    p.paragraph_format.space_before = Pt(0)
    run = p.add_run(text)
    set_font(run, name=font, size=size, bold=bold, color=color)
    if fill is not None:
        shade_cell(cell, fill)
    set_cell_borders(cell, LIGHT, 0.5)


def insert_logo(container, width_in=2.0):
    if hasattr(container, "add_picture"):
        return container.add_picture(str(LOGO_PNG), width=Inches(width_in))
    p = container.paragraphs[0] if not container.text else container.add_paragraph()
    run = p.add_run()
    run.add_picture(str(LOGO_PNG), width=Inches(width_in))
    return p


def add_disclaimer(section_or_doc, text=DISCLAIMER):
    return add_para(section_or_doc, text, font=BODY_FONT, size=8, color=MID, space_after=0)


def base_doc(margin_in=0.8):
    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(margin_in)
        section.bottom_margin = Inches(margin_in)
        section.left_margin = Inches(margin_in)
        section.right_margin = Inches(margin_in)
    style = doc.styles["Normal"]
    style.font.name = BODY_FONT
    style.font.size = Pt(11)
    style.font.color.rgb = CHARCOAL
    return doc


def add_letterhead_header(doc, address_lines=None):
    section = doc.sections[0]
    header = section.header
    htable = header.add_table(rows=1, cols=2, width=Inches(7))
    htable.autofit = False
    htable.columns[0].width = Inches(3)
    htable.columns[1].width = Inches(4)
    htable.rows[0].height = Inches(0.9)

    left = htable.cell(0, 0)
    left.paragraphs[0].text = ""
    insert_logo(left, width_in=2.0)

    right = htable.cell(0, 1)
    right.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    right.paragraphs[0].text = ""
    addr = address_lines or [
        "VisionWise Capital",
        "[Street Address]",
        "[City, State ZIP]",
        "[Phone] · visionwisecapital.com",
    ]
    for i, line in enumerate(addr):
        p = right.paragraphs[0] if i == 0 else right.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.paragraph_format.space_after = Pt(0)
        run = p.add_run(line)
        set_font(run, size=9, color=MID if i > 0 else CHARCOAL, bold=(i == 0))


def add_letterhead_footer(doc, with_disclaimer=True):
    section = doc.sections[0]
    footer = section.footer
    p1 = footer.paragraphs[0]
    p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p1.add_run("VisionWise Capital  ·  visionwisecapital.com  ·  [Phone]  ·  [Email]")
    set_font(run, size=9, color=MID)
    if with_disclaimer:
        p2 = footer.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run2 = p2.add_run(DISCLAIMER)
        set_font(run2, size=7, color=MID)


# ---------------------------------------------------------------------------
# DOCX templates
# ---------------------------------------------------------------------------

def build_letterhead():
    doc = base_doc(0.8)
    add_letterhead_header(doc)
    add_letterhead_footer(doc, with_disclaimer=False)

    add_para(doc, "[Date]", size=11, color=CHARCOAL, space_after=18)
    add_para(doc, "[Recipient Name]", size=11, color=CHARCOAL, space_after=0)
    add_para(doc, "[Title]", size=11, color=CHARCOAL, space_after=0)
    add_para(doc, "[Company]", size=11, color=CHARCOAL, space_after=0)
    add_para(doc, "[Address]", size=11, color=CHARCOAL, space_after=18)

    add_para(doc, "Dear [First Name],", size=11, color=CHARCOAL, space_after=12)
    add_para(doc,
             "[Body paragraph 1 — open warmly. Reference the prior conversation or shared connection.]",
             size=11, space_after=10)
    add_para(doc,
             "[Body paragraph 2 — make the specific point or offer. Keep it concrete and short.]",
             size=11, space_after=10)
    add_para(doc,
             "[Body paragraph 3 — close with a clear next step (call time, doc to review, intro to make).]",
             size=11, space_after=24)

    add_para(doc, "Sincerely,", size=11, space_after=36)
    add_para(doc, "Sanford D. Coggins", size=11, bold=True, color=CHARCOAL, space_after=0)
    add_para(doc, "Founder & CEO, VisionWise Capital", size=10, color=MID, space_after=0)

    out = TEMPLATES / "letterhead.docx"
    doc.save(out)
    return out


def build_investor_letter():
    doc = base_doc(0.8)
    add_letterhead_header(doc)
    add_letterhead_footer(doc, with_disclaimer=True)

    add_para(doc, "[Date]", size=11, space_after=18)
    add_para(doc, "[Investor Name]", size=11, space_after=0)
    add_para(doc, "[Address]", size=11, space_after=18)

    add_para(doc, "Re: [Deal Name] — Updated Offering Materials", size=11, bold=True, color=GREEN_DARK, space_after=12)

    add_para(doc, "Dear [First Name],", size=11, space_after=12)
    add_para(doc,
             "Thank you for the time we spent discussing [Deal Name] last week. As promised, "
             "I'm sending the updated materials so you can review them at your own pace.",
             space_after=10)
    add_para(doc, "Enclosed you'll find:", space_after=4)
    for item in [
        "Offering one-pager and full Private Placement Memorandum",
        "Property condition report and third-party appraisal",
        "Pro forma model with sensitivity analysis",
        "VisionWise transaction history (full track record)",
    ]:
        add_para(doc, "  •  " + item, size=11, space_after=2)
    add_para(doc, "", space_after=8)
    add_para(doc,
             "I want to be candid about how we're thinking about risk on this deal — see "
             "page [X] of the memorandum for the specific assumptions we're underwriting "
             "against. Happy to walk through any line of the model on a call.",
             space_after=10)
    add_para(doc,
             "Next step: I'll follow up Friday to answer questions. If you'd like to move "
             "forward, the soft-commit window closes [Date].",
             space_after=24)

    add_para(doc, "With gratitude,", space_after=36)
    add_para(doc, "Sanford D. Coggins", size=11, bold=True, space_after=0)
    add_para(doc, "Founder & CEO, VisionWise Capital", size=10, color=MID, space_after=0)

    out = TEMPLATES / "investor-letter.docx"
    doc.save(out)
    return out


def build_quarterly_report():
    doc = base_doc(0.7)

    # Cover band
    table = doc.add_table(rows=1, cols=1)
    table.autofit = False
    table.columns[0].width = Inches(7.0)
    cell = table.cell(0, 0)
    shade_cell(cell, CHARCOAL)
    cell.paragraphs[0].text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(36)
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run("Quarterly Investor Report")
    set_font(run, name=HEADING_FONT, size=11, color=GREEN, bold=True)
    p2 = cell.add_paragraph()
    p2.paragraph_format.space_after = Pt(0)
    run2 = p2.add_run("[Q# YYYY]")
    set_font(run2, name=HEADING_FONT, size=36, color=WHITE, bold=True)
    p3 = cell.add_paragraph()
    p3.paragraph_format.space_after = Pt(36)
    run3 = p3.add_run("VisionWise Capital  ·  Portfolio Update")
    set_font(run3, name=BODY_FONT, size=11, color=WHITE)

    add_para(doc, "", space_after=6)
    insert_logo(doc, width_in=1.6)

    add_para(doc, "", space_after=12)
    h2(doc, "A note from Sanford")
    add_para(doc,
             "[2–4 paragraph letter. Cover the quarter honestly: what worked, what didn't, "
             "what changed in our view of the market, and what we're paying attention to next "
             "quarter. Keep it plain and personal — investors are reading this for your "
             "judgment, not the numbers (those are below).]",
             space_after=24)

    h2(doc, "Portfolio at a glance")
    metrics = doc.add_table(rows=2, cols=4)
    metrics.autofit = False
    for col in metrics.columns:
        col.width = Inches(1.75)
    headers = ["Active deals", "Total AUM", "Distributions YTD", "Avg. cash-on-cash"]
    values = ["[#]", "$[##.#]M", "$[###]K", "[#.#]%"]
    for i, h in enumerate(headers):
        fill_cell(metrics.cell(0, i), h, size=9, bold=True, color=MID, fill=OFFWHITE, align=WD_ALIGN_PARAGRAPH.CENTER)
    for i, v in enumerate(values):
        fill_cell(metrics.cell(1, i), v, size=18, bold=True, color=GREEN_DARK, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, "", space_after=14)

    h2(doc, "Per-deal updates")
    deal_table = doc.add_table(rows=4, cols=5)
    headers = ["Deal", "Status", "Hold to date", "Distributions QTD", "Notes"]
    for i, htext in enumerate(headers):
        fill_cell(deal_table.cell(0, i), htext, size=9, bold=True, color=WHITE, fill=CHARCOAL,
                  align=WD_ALIGN_PARAGRAPH.LEFT)
    sample_rows = [
        ("[Property Name]", "[On plan]", "[##] mo", "$[##]K", "[1-line update]"),
        ("[Property Name]", "[Refinance]", "[##] mo", "$[##]K", "[1-line update]"),
        ("[Property Name]", "[Lease-up]", "[##] mo", "$[##]K", "[1-line update]"),
    ]
    for r, row in enumerate(sample_rows, start=1):
        for c, val in enumerate(row):
            fill_cell(deal_table.cell(r, c), val, size=10, color=CHARCOAL, align=WD_ALIGN_PARAGRAPH.LEFT)
    add_para(doc, "", space_after=14)

    h2(doc, "Pipeline & outlook")
    add_para(doc,
             "[2–3 paragraphs on what's under contract, what's in diligence, and the macro "
             "or local-market shifts you're tracking. End with what investors should expect "
             "next quarter and any capital calls.]",
             space_after=18)

    h2(doc, "Important disclosures")
    add_para(doc, DISCLAIMER, size=8, color=MID, space_after=0)

    out = TEMPLATES / "investor-quarterly-report.docx"
    doc.save(out)
    return out


def build_case_study():
    doc = base_doc(0.7)

    # Top label
    add_para(doc, "CASE STUDY  ·  CLOSED DEAL", size=10, bold=True, color=GREEN_DARK, space_after=2)
    h1(doc, "[Property Name]")
    add_para(doc, "[City, State]  ·  [Asset class]  ·  [# units]", size=11, color=MID, space_after=14)

    # Hero photo placeholder
    photo_table = doc.add_table(rows=1, cols=1)
    photo_table.autofit = False
    photo_table.columns[0].width = Inches(7.0)
    cell = photo_table.cell(0, 0)
    cell.height = Inches(2.6)
    shade_cell(cell, OFFWHITE)
    set_cell_borders(cell, LIGHT, 0.5)
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(80)
    run = p.add_run("[ Insert hero photo — before/after side-by-side preferred ]")
    set_font(run, size=10, color=MID, italic=True)
    add_para(doc, "", space_after=14)

    # Snapshot table
    h2(doc, "Snapshot")
    snap = doc.add_table(rows=4, cols=4)
    snap.autofit = False
    rows = [
        ("Acquired", "[MM/YYYY]", "Sold / refi'd", "[MM/YYYY]"),
        ("Purchase price", "$[##.#]M", "Exit / value", "$[##.#]M"),
        ("Hold period", "[##] months", "Total equity", "$[##.#]M"),
        ("Investor IRR", "[##.#]%", "Equity multiple", "[#.##]x"),
    ]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            bold = (c % 2 == 0)
            color = MID if bold else CHARCOAL
            fill = OFFWHITE if bold else None
            fill_cell(snap.cell(r, c), val, size=10, bold=bold, color=color, fill=fill,
                      align=WD_ALIGN_PARAGRAPH.LEFT)
    add_para(doc, "", space_after=14)

    h2(doc, "The thesis")
    add_para(doc,
             "[Why we bought it. 2–3 sentences. What was mispriced, mismanaged, or "
             "undervalued — and why we believed our team could fix it.]",
             space_after=14)

    h2(doc, "What we did")
    for bullet in [
        "[Operational change — e.g., re-tenanted, repositioned units, stabilized expenses.]",
        "[Capital improvement — e.g., renovated kitchens, exterior, common areas.]",
        "[Financial structure — e.g., refinanced at month [##], distributed [##]% of equity.]",
    ]:
        add_para(doc, "  •  " + bullet, space_after=4)
    add_para(doc, "", space_after=14)

    h2(doc, "The result")
    add_para(doc,
             "[Plain-language summary of returns and what investors actually received. "
             "Reference the equity multiple and IRR in dollar terms — \"$100K invested became "
             "$[##]K\" — so a non-finance reader gets it instantly.]",
             space_after=14)

    h2(doc, "Investor reflection")
    add_para(doc,
             "“[Pull a real, named-or-anonymized investor quote. Keep it short and "
             "specific.]”",
             italic=True, color=CHARCOAL, space_after=4)
    add_para(doc, "— [Investor Name or “LP since 2018”]", size=10, color=MID, space_after=18)

    add_para(doc, DISCLAIMER, size=8, color=MID, space_after=0)

    out = TEMPLATES / "case-study.docx"
    doc.save(out)
    return out


def build_offering_one_pager():
    doc = base_doc(0.6)

    # Header band with logo
    band = doc.add_table(rows=1, cols=2)
    band.autofit = False
    band.columns[0].width = Inches(2.4)
    band.columns[1].width = Inches(5.0)
    insert_logo(band.cell(0, 0), width_in=1.8)
    right = band.cell(0, 1)
    right.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    p = right.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run("INVESTMENT OPPORTUNITY")
    set_font(run, size=10, bold=True, color=GREEN_DARK)
    add_para(doc, "", space_after=6)

    h1(doc, "[Property Name]")
    add_para(doc, "[Asset class]  ·  [City, State]  ·  [# units]", size=12, color=MID, space_after=10)

    # Hero placeholder
    hero = doc.add_table(rows=1, cols=1)
    hero.columns[0].width = Inches(7.4)
    hcell = hero.cell(0, 0)
    shade_cell(hcell, OFFWHITE)
    set_cell_borders(hcell, LIGHT, 0.5)
    php = hcell.paragraphs[0]
    php.alignment = WD_ALIGN_PARAGRAPH.CENTER
    php.paragraph_format.space_before = Pt(60)
    php.paragraph_format.space_after = Pt(60)
    run = php.add_run("[ Insert property photo ]")
    set_font(run, size=10, color=MID, italic=True)
    add_para(doc, "", space_after=12)

    # 4-up KPI strip
    kpis = doc.add_table(rows=2, cols=4)
    kpi_labels = ["Target IRR", "Avg. Cash-on-Cash", "Hold Period", "Min. Investment"]
    kpi_values = ["[##]%", "[#.#]%", "[#] yrs", "$[###]K"]
    for i, l in enumerate(kpi_labels):
        fill_cell(kpis.cell(0, i), l, size=9, bold=True, color=MID, fill=OFFWHITE,
                  align=WD_ALIGN_PARAGRAPH.CENTER)
    for i, v in enumerate(kpi_values):
        fill_cell(kpis.cell(1, i), v, size=18, bold=True, color=GREEN_DARK,
                  align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, "", space_after=14)

    h2(doc, "Why this deal")
    for bullet in [
        "[Thesis bullet 1 — the mispricing or value-add lever.]",
        "[Thesis bullet 2 — the operational plan and timeline.]",
        "[Thesis bullet 3 — the downside protection.]",
    ]:
        add_para(doc, "  •  " + bullet, space_after=4)
    add_para(doc, "", space_after=10)

    h2(doc, "Why VisionWise")
    add_para(doc,
             "Since [year], VisionWise Capital has acquired and operated [##] properties "
             "in Southern California with [##.#]M in investor equity deployed. We invest "
             "alongside our LPs on every deal.",
             space_after=14)

    h2(doc, "Next step")
    add_para(doc,
             "Schedule a consultation with Sanford to walk through the model and "
             "diligence package.  ·  visionwisecapital.com  ·  [Phone]  ·  [Email]",
             bold=True, color=GREEN_DARK, space_after=14)

    add_para(doc, DISCLAIMER, size=8, color=MID, space_after=0)

    out = TEMPLATES / "offering-one-pager.docx"
    doc.save(out)
    return out


def build_press_release():
    doc = base_doc(0.9)
    add_letterhead_footer(doc, with_disclaimer=False)

    insert_logo(doc, width_in=1.8)
    add_para(doc, "", space_after=6)

    add_para(doc, "FOR IMMEDIATE RELEASE", size=10, bold=True, color=GREEN_DARK, space_after=14)

    h1(doc, "[Headline — One Sharp Line, Sentence Case]")
    add_para(doc, "[Subhead — clarifying line in plain language, one sentence.]",
             font=HEADING_FONT, size=14, color=MID, space_after=14)

    add_para(doc,
             "[CITY, STATE — Month DD, YYYY] — VisionWise Capital, a Southern "
             "California-based private real estate investment firm, today announced "
             "[the news, in one sentence].",
             space_after=10)
    add_para(doc,
             "[Paragraph 2 — supporting detail. Numbers, dates, and what this means for "
             "investors or the local market.]",
             space_after=10)
    add_para(doc,
             "“[Quote from Sanford D. Coggins, 1–2 sentences, in his actual voice — "
             "stewardship, long-term, plain English.]” said Sanford D. Coggins, "
             "Founder and CEO of VisionWise Capital.",
             space_after=10)
    add_para(doc,
             "[Paragraph 4 — context, history, or what's next.]",
             space_after=18)

    h3(doc, "About VisionWise Capital")
    add_para(doc,
             "VisionWise Capital is a private real estate investment firm based in "
             "Southern California. The firm acquires and operates multifamily and "
             "residential properties on behalf of accredited investors, with a focus on "
             "long-term value, transparency, and disciplined underwriting. Learn more at "
             "visionwisecapital.com.",
             space_after=14)

    h3(doc, "Media contact")
    add_para(doc, "[Name]", space_after=0)
    add_para(doc, "[Title]", size=10, color=MID, space_after=0)
    add_para(doc, "[Email]  ·  [Phone]", size=10, color=MID, space_after=10)

    add_para(doc, "###", size=11, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, color=MID)

    out = TEMPLATES / "press-release.docx"
    doc.save(out)
    return out


def build_memo():
    doc = base_doc(0.8)

    insert_logo(doc, width_in=1.4)
    add_para(doc, "", space_after=6)
    h1(doc, "Internal Memorandum")
    p = doc.paragraphs[-1]
    add_bottom_border(p, GREEN, size_pt=2)

    meta = doc.add_table(rows=4, cols=2)
    meta.autofit = False
    meta.columns[0].width = Inches(0.9)
    meta.columns[1].width = Inches(6.5)
    rows = [("To:", "[Name(s)]"), ("From:", "[Name]"), ("Date:", "[Date]"), ("Re:", "[Subject — short and concrete]")]
    for i, (label, val) in enumerate(rows):
        fill_cell(meta.cell(i, 0), label, size=10, bold=True, color=MID, align=WD_ALIGN_PARAGRAPH.LEFT)
        fill_cell(meta.cell(i, 1), val, size=11, color=CHARCOAL, align=WD_ALIGN_PARAGRAPH.LEFT)
        for cc in (meta.cell(i, 0), meta.cell(i, 1)):
            tcPr = cc._tc.get_or_add_tcPr()
            for b in tcPr.findall(qn("w:tcBorders")):
                tcPr.remove(b)
    add_para(doc, "", space_after=10)

    h2(doc, "Summary")
    add_para(doc, "[2–3 sentences. The decision or update, up front.]", space_after=14)

    h2(doc, "Context")
    add_para(doc, "[What's the situation. What changed. Who's affected.]", space_after=14)

    h2(doc, "Recommendation / decision")
    add_para(doc, "[The actual call, plus the reasoning in plain language.]", space_after=14)

    h2(doc, "Action items")
    actions = doc.add_table(rows=4, cols=3)
    headers = ["Owner", "Action", "Due"]
    for i, htxt in enumerate(headers):
        fill_cell(actions.cell(0, i), htxt, size=9, bold=True, color=WHITE, fill=CHARCOAL,
                  align=WD_ALIGN_PARAGRAPH.LEFT)
    for r in range(1, 4):
        fill_cell(actions.cell(r, 0), "[Name]", size=10, align=WD_ALIGN_PARAGRAPH.LEFT)
        fill_cell(actions.cell(r, 1), "[What]", size=10, align=WD_ALIGN_PARAGRAPH.LEFT)
        fill_cell(actions.cell(r, 2), "[Date]", size=10, align=WD_ALIGN_PARAGRAPH.LEFT)

    out = TEMPLATES / "memo-internal-update.docx"
    doc.save(out)
    return out


# ---------------------------------------------------------------------------
# Email signature
# ---------------------------------------------------------------------------

EMAIL_SIG = """<!doctype html>
<html><body style="margin:0;padding:0;">
<!--
  VisionWise Capital — HTML email signature.
  Open this file, copy the rendered signature into Outlook/Gmail, then replace
  bracketed values. Inline styles only — required for email-client compat.

  Font stack is brand-first with Google Fonts fallback. Most email clients
  will substitute system sans (Helvetica/Arial) regardless — that's expected.
-->
<table cellpadding="0" cellspacing="0" border="0" style="font-family:'Proxima Nova','Open Sans',Helvetica,Arial,sans-serif;color:#40464B;font-size:13px;line-height:1.45;">
  <tr>
    <td style="padding-right:18px;border-right:3px solid #A3D55D;vertical-align:top;">
      <img src="https://visionwisecapital.com/wp-content/uploads/logo-small.png"
           alt="VisionWise Capital"
           width="180" height="29"
           style="display:block;border:0;">
    </td>
    <td style="padding-left:18px;vertical-align:top;">
      <div style="font-family:'Gill Sans MT Pro','Lato',Helvetica,Arial,sans-serif;font-size:15px;font-weight:700;color:#40464B;">[Full Name]</div>
      <div style="font-size:12px;color:#8A8F94;">[Title]  ·  VisionWise Capital</div>
      <div style="height:6px;line-height:6px;">&nbsp;</div>
      <div><a href="tel:+1[##########]" style="color:#40464B;text-decoration:none;">[Phone]</a>
        &nbsp;·&nbsp;
        <a href="mailto:[email]@visionwisecapital.com" style="color:#5C9627;text-decoration:none;">[email]@visionwisecapital.com</a>
      </div>
      <div><a href="https://visionwisecapital.com" style="color:#5C9627;text-decoration:none;">visionwisecapital.com</a></div>
    </td>
  </tr>
  <tr>
    <td colspan="2" style="padding-top:12px;font-size:9px;color:#8A8F94;line-height:1.4;font-family:'Proxima Nova','Open Sans',Helvetica,Arial,sans-serif;">
      This email and any attachments are confidential. Nothing in this message constitutes an
      offer to sell or a solicitation of an offer to buy securities. Real estate investments
      involve risk, including loss of principal.
    </td>
  </tr>
</table>
</body></html>
"""


def build_email_signature():
    out = TEMPLATES / "email-signature.html"
    out.write_text(EMAIL_SIG, encoding="utf-8")
    return out


# ---------------------------------------------------------------------------
# PPTX deck
# ---------------------------------------------------------------------------

def add_pptx_text(slide, left, top, width, height, text, *, font=BODY_FONT, size=18,
                  bold=False, color=P_CHARCOAL, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = PPt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return box


def add_color_block(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_logo_to_slide(slide, left, top, width):
    return slide.shapes.add_picture(str(LOGO_PNG), left, top, width=width)


def build_pitch_deck():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height

    # 1. Title slide — charcoal background, green accent strip
    s = prs.slides.add_slide(blank)
    add_color_block(s, 0, 0, SLIDE_W, SLIDE_H, P_CHARCOAL)
    add_color_block(s, 0, 0, PInches(0.35), SLIDE_H, P_GREEN)
    add_logo_to_slide(s, PInches(0.9), PInches(0.7), PInches(2.4))
    add_pptx_text(s, PInches(0.9), PInches(3.4), PInches(11), PInches(0.5),
                  "INVESTMENT OPPORTUNITY", font=HEADING_FONT, size=14, bold=True, color=P_GREEN)
    add_pptx_text(s, PInches(0.9), PInches(3.9), PInches(11), PInches(1.6),
                  "[Deal Name]", font=HEADING_FONT, size=54, bold=True, color=P_WHITE)
    add_pptx_text(s, PInches(0.9), PInches(5.4), PInches(11), PInches(0.5),
                  "[City, State]  ·  [# Units]  ·  [Asset Class]",
                  font=BODY_FONT, size=18, color=P_GREEN_TINT)
    add_pptx_text(s, PInches(0.9), PInches(6.7), PInches(11), PInches(0.4),
                  "VisionWise Capital  ·  [Date]  ·  Confidential — Accredited Investors Only",
                  font=BODY_FONT, size=10, color=P_GREEN_TINT)

    # 2. Disclaimer / safe harbor
    s = prs.slides.add_slide(blank)
    add_color_block(s, 0, 0, SLIDE_W, PInches(0.35), P_GREEN)
    add_pptx_text(s, PInches(0.6), PInches(0.7), PInches(12), PInches(0.6),
                  "Important Disclosures", font=HEADING_FONT, size=28, bold=True, color=P_CHARCOAL)
    add_pptx_text(s, PInches(0.6), PInches(1.5), PInches(12), PInches(5.5),
                  DISCLAIMER + "\n\nThis presentation contains forward-looking statements based "
                  "on assumptions and analyses made by VisionWise Capital that we believe are "
                  "reasonable, but actual results may differ materially. Recipients should "
                  "consult their own tax, legal, and financial advisors before investing.",
                  font=BODY_FONT, size=14, color=P_CHARCOAL)

    # 3. Section divider helper
    def section(title, kicker=None):
        s = prs.slides.add_slide(blank)
        add_color_block(s, 0, 0, SLIDE_W, SLIDE_H, P_CHARCOAL)
        add_color_block(s, 0, PInches(3.4), PInches(1.2), PInches(0.08), P_GREEN)
        if kicker:
            add_pptx_text(s, PInches(0.9), PInches(3.0), PInches(12), PInches(0.4),
                          kicker.upper(), font=HEADING_FONT, size=12, bold=True, color=P_GREEN)
        add_pptx_text(s, PInches(0.9), PInches(3.55), PInches(12), PInches(1.2),
                      title, font=HEADING_FONT, size=44, bold=True, color=P_WHITE)
        return s

    # Content slide helper
    def content(title):
        s = prs.slides.add_slide(blank)
        add_color_block(s, 0, 0, SLIDE_W, PInches(0.35), P_GREEN)
        add_logo_to_slide(s, PInches(11.6), PInches(0.55), PInches(1.5))
        add_pptx_text(s, PInches(0.6), PInches(0.6), PInches(11), PInches(0.7),
                      title, font=HEADING_FONT, size=28, bold=True, color=P_CHARCOAL)
        add_color_block(s, PInches(0.6), PInches(1.35), PInches(0.8), PInches(0.04), P_GREEN)
        add_pptx_text(s, PInches(0.6), PInches(7.05), PInches(12), PInches(0.3),
                      "VisionWise Capital  ·  Confidential",
                      font=BODY_FONT, size=9, color=P_MID)
        return s

    section("VisionWise Capital", kicker="Who we are")
    s = content("About VisionWise")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                  "[2–3 plain sentences on what VWC does, who we serve, and what makes our "
                  "discipline different. End with the core promise: stewardship, "
                  "long-term thinking, relationship-first.]\n\n"
                  "•  Founded by Sanford D. Coggins\n"
                  "•  Southern California focus, multifamily + residential\n"
                  "•  We invest alongside every LP on every deal",
                  font=BODY_FONT, size=18, color=P_CHARCOAL)

    s = content("Track record")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(0.5),
                  "[High-level track record line — # deals, total equity, average IRR.]",
                  font=BODY_FONT, size=18, color=P_CHARCOAL)
    # 4-up KPI cards
    kpi_y = PInches(2.6)
    kpi_w = PInches(2.85)
    kpi_h = PInches(1.6)
    kpi_x_start = PInches(0.6)
    gap = PInches(0.2)
    labels = ["Deals closed", "Total equity", "Avg. Investor IRR", "Avg. EM"]
    values = ["[##]", "$[##]M", "[##]%", "[#.##]x"]
    for i, (label, val) in enumerate(zip(labels, values)):
        x = kpi_x_start + i * (kpi_w + gap)
        add_color_block(s, x, kpi_y, kpi_w, kpi_h, P_GREEN_TINT)
        add_pptx_text(s, x + PInches(0.2), kpi_y + PInches(0.2), kpi_w - PInches(0.4), PInches(0.3),
                      label.upper(), font=BODY_FONT, size=10, bold=True, color=P_GREEN_DARK)
        add_pptx_text(s, x + PInches(0.2), kpi_y + PInches(0.55), kpi_w - PInches(0.4), PInches(1.0),
                      val, font=HEADING_FONT, size=36, bold=True, color=P_CHARCOAL)
    add_pptx_text(s, PInches(0.6), PInches(4.5), PInches(12), PInches(2.3),
                  "[Bullet list of 3–5 representative closed deals with one-line outcome each.]",
                  font=BODY_FONT, size=16, color=P_CHARCOAL)

    section("The Opportunity", kicker="The deal")
    s = content("Property overview")
    add_color_block(s, PInches(0.6), PInches(1.7), PInches(6.0), PInches(4.5), P_LIGHT)
    add_pptx_text(s, PInches(0.6), PInches(3.7), PInches(6.0), PInches(0.4),
                  "[ Insert hero photo ]", font=BODY_FONT, size=12, italic=True, color=P_MID, align=PP_ALIGN.CENTER)
    add_pptx_text(s, PInches(7.0), PInches(1.7), PInches(5.8), PInches(4.8),
                  "[Property name]\n[Address, City, State]\n\n"
                  "•  [Asset class / # units]\n"
                  "•  [Year built / Year renovated]\n"
                  "•  [Sq ft / Lot size]\n"
                  "•  [Current occupancy]\n"
                  "•  [In-place vs market rent gap]\n",
                  font=BODY_FONT, size=16, color=P_CHARCOAL)

    s = content("Investment thesis")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                  "•  [Why this is mispriced or under-managed today]\n\n"
                  "•  [The specific value-creation plan and timeline]\n\n"
                  "•  [What protects the downside — debt structure, market depth, sponsor experience]",
                  font=BODY_FONT, size=20, color=P_CHARCOAL)

    s = content("Market & comps")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.2),
                  "[Market context: rent trends, supply pipeline, demographics, employment. "
                  "End with a clear sentence on why THIS submarket THIS year.]\n\n"
                  "[Comp table or chart placeholder — replace with actual rent comps and sale comps.]",
                  font=BODY_FONT, size=16, color=P_CHARCOAL)

    section("Returns & Structure", kicker="The numbers")
    s = content("Projected returns")
    headers = ["", "Base case", "Downside", "Upside"]
    rows_lines = [
        ("Equity multiple", "[#.##]x", "[#.##]x", "[#.##]x"),
        ("Investor IRR", "[##]%", "[##]%", "[##]%"),
        ("Avg. cash-on-cash", "[#.#]%", "[#.#]%", "[#.#]%"),
        ("Hold period", "[#] yrs", "[#] yrs", "[#] yrs"),
    ]
    table_x, table_y = PInches(0.6), PInches(1.8)
    col_w = [PInches(3.6), PInches(2.7), PInches(2.7), PInches(2.7)]
    row_h = PInches(0.55)
    # header
    cx = table_x
    for i, htxt in enumerate(headers):
        add_color_block(s, cx, table_y, col_w[i], row_h, P_CHARCOAL)
        add_pptx_text(s, cx + PInches(0.15), table_y + PInches(0.12), col_w[i] - PInches(0.3), PInches(0.4),
                      htxt, font=BODY_FONT, size=14, bold=True, color=P_WHITE)
        cx += col_w[i]
    # rows
    for ri, row in enumerate(rows_lines):
        ry = table_y + row_h * (ri + 1)
        if ri % 2 == 0:
            add_color_block(s, table_x, ry, sum(col_w, Emu(0)), row_h, PRGB(0xF7, 0xF7, 0xF7))
        cx = table_x
        for i, val in enumerate(row):
            add_pptx_text(s, cx + PInches(0.15), ry + PInches(0.13), col_w[i] - PInches(0.3), PInches(0.4),
                          val, font=BODY_FONT, size=14, bold=(i == 0), color=P_CHARCOAL)
            cx += col_w[i]

    s = content("Deal structure & fees")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                  "•  Total raise:  $[##]M\n"
                  "•  Min. investment:  $[###]K\n"
                  "•  Preferred return:  [#]% to LPs\n"
                  "•  Promote / waterfall:  [##/##] over the pref\n"
                  "•  GP co-invest:  [##]% (we invest alongside you)\n"
                  "•  Asset management fee:  [#.##]% of equity\n"
                  "•  Acquisition fee:  [#.##]% of purchase price",
                  font=BODY_FONT, size=18, color=P_CHARCOAL)

    section("Why VisionWise", kicker="Our edge")
    s = content("Our team")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.2),
                  "[Brief on Sanford D. Coggins and the operating team. Years in market, "
                  "deals executed, what they bring to this specific deal. 2–3 sentences each.]",
                  font=BODY_FONT, size=18, color=P_CHARCOAL)

    s = content("Risks we're underwriting")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                  "•  [Market risk — and how we mitigate it]\n\n"
                  "•  [Execution risk — and the team's experience with this specific lift]\n\n"
                  "•  [Capital structure risk — and the debt terms]\n\n"
                  "•  [Liquidity — and the expected hold]",
                  font=BODY_FONT, size=18, color=P_CHARCOAL)

    section("Next Steps", kicker="How to invest")
    s = content("Process & timeline")
    add_pptx_text(s, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                  "1.  Soft-commit  —  by [Date]\n\n"
                  "2.  Diligence call & document review  —  weeks of [Date]\n\n"
                  "3.  Subscription agreement & wire  —  by [Date]\n\n"
                  "4.  Close  —  [Target close date]",
                  font=BODY_FONT, size=20, color=P_CHARCOAL)

    s = prs.slides.add_slide(blank)
    add_color_block(s, 0, 0, SLIDE_W, SLIDE_H, P_CHARCOAL)
    add_color_block(s, 0, 0, PInches(0.35), SLIDE_H, P_GREEN)
    add_logo_to_slide(s, PInches(0.9), PInches(0.8), PInches(2.4))
    add_pptx_text(s, PInches(0.9), PInches(3.1), PInches(11), PInches(0.5),
                  "LET'S TALK", font=HEADING_FONT, size=14, bold=True, color=P_GREEN)
    add_pptx_text(s, PInches(0.9), PInches(3.6), PInches(11), PInches(1.2),
                  "Sanford D. Coggins", font=HEADING_FONT, size=44, bold=True, color=P_WHITE)
    add_pptx_text(s, PInches(0.9), PInches(4.7), PInches(11), PInches(0.5),
                  "Founder & CEO, VisionWise Capital", font=BODY_FONT, size=18, color=P_GREEN_TINT)
    add_pptx_text(s, PInches(0.9), PInches(5.6), PInches(11), PInches(0.5),
                  "[email]@visionwisecapital.com  ·  [Phone]", font=BODY_FONT, size=16, color=P_WHITE)
    add_pptx_text(s, PInches(0.9), PInches(6.1), PInches(11), PInches(0.5),
                  "visionwisecapital.com", font=BODY_FONT, size=16, color=P_GREEN)

    out = TEMPLATES / "pitch-deck.pptx"
    prs.save(out)
    return out


# ---------------------------------------------------------------------------

def main():
    built = []
    for fn in (
        build_letterhead,
        build_investor_letter,
        build_quarterly_report,
        build_case_study,
        build_offering_one_pager,
        build_press_release,
        build_memo,
        build_email_signature,
        build_pitch_deck,
    ):
        path = fn()
        built.append(path)
        print(f"  built  {path.relative_to(ROOT)}")
    print(f"\n{len(built)} templates written to {TEMPLATES.relative_to(ROOT)}/")


if __name__ == "__main__":
    main()
